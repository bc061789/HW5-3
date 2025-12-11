import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO

st.set_page_config(page_title="AI PPT Re-Designer", page_icon="🧠")
st.title("🧠 AI PowerPoint 版型重新設計 Demo")

st.markdown("""
這個 demo 示範：  
1. 上傳一份原始 PPTX  
2. 選擇一種「AI 設計風格」  
3. 自動產生一份 **重新排版的 PPT** 並提供下載  
""")

uploaded = st.file_uploader("請上傳一份 PPTX 檔案", type=["pptx"])

style = st.radio(
    "選擇重新設計風格",
    ["科技藍 Tech Style", "極簡白 Minimal Style"]
)

if uploaded is not None:
    st.info("檔案已上傳：共 {} 頁投影片".format(len(Presentation(uploaded).slides)))

if uploaded and st.button("🚀 產生新的 PPT"):
    # 讀取原始簡報
    old_prs = Presentation(uploaded)

    # 建立新的簡報（先清空預設投影片）
    new_prs = Presentation()
    while len(new_prs.slides) > 0:
        r_id = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(r_id)
        del new_prs.slides._sldIdLst[0]

    for old_slide in old_prs.slides:
        layout = new_prs.slide_layouts[6]  # blank
        slide = new_prs.slides.add_slide(layout)

        # 設定背景顏色
        bg_fill = slide.background.fill
        bg_fill.solid()
        if style.startswith("科技藍"):
            bg_fill.fore_color.rgb = RGBColor(8, 24, 72)   # 深藍
            font_color = RGBColor(255, 255, 255)           # 白字
        else:
            bg_fill.fore_color.rgb = RGBColor(255, 255, 255)  # 全白
            font_color = RGBColor(40, 40, 40)                 # 深灰字

        # 把舊投影片中的文字，簡單「抽出來」重排成一列一列的文字框
        top = Inches(1)
        for shape in old_slide.shapes:
            if not shape.has_text_frame:
                continue

            textbox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(1))
            tf = textbox.text_frame
            tf.text = shape.text

            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(24)
                    r.font.color.rgb = font_color

            top += Inches(0.8)

        # 極簡版多加一條細線裝飾
        if style.startswith("極簡白"):
            line = slide.shapes.add_shape(
                autoshape_type_id=1,  # rectangle
                left=Inches(0.8),
                top=Inches(0.8),
                width=Inches(0.05),
                height=Inches(6)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(210, 180, 90)  # 淺金色
            line.line.fill.background()

    # 輸出到記憶體並提供下載
    output = BytesIO()
    new_prs.save(output)
    output.seek(0)

    st.success("✅ 重新設計完成！請下載新的 PPT 檔案。")

    filename = "redesigned_tech_style.pptx" if style.startswith("科技藍") else "redesigned_minimal_style.pptx"
    st.download_button(
        label="💾 下載新 PPT",
        data=output,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
else:
    st.caption("⬆ 先上傳一份 PPT，才能產生新的檔案。")
